from typing import Iterator
from langchain_core.documents import Document
from langchain_core.document_loaders import BaseLoader
from pptx import Presentation

class CustomPPTLoader(BaseLoader):
    def __init__(self, file_path: str) -> None:
        self.file_path = file_path

    def lazy_load(self) -> Iterator[Document]:
        prs = Presentation(self.file_path)
        for slide_num, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            texts.append(paragraph.text)
            if texts:
                yield Document(
                    page_content="\n".join(texts),
                    metadata={"source": self.file_path, "slide": slide_num}
                )

if __name__ == '__main__':
    import os
    absfile = os.path.abspath(__file__)
    dirfile = os.path.dirname(absfile)
    dirfile = os.path.dirname(dirfile)
    file = os.path.join(dirfile, 'data/test_ppt.pptx')
    loader = CustomPPTLoader(file)  # 注意：笔记中定义的类名是 CustomPPTLoader，项目实际代码可能用 OCRPPTLoader
    docs = loader.load()
    for doc in docs:
        print(f'幻灯片 {doc.metadata.get("slide")}: {doc.page_content[:80]}')
